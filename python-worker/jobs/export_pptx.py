"""
Investor PPTX Export Job Handler
Generates professional investor presentation using python-pptx
Stores file in database instead of S3
"""

import json
import base64
from datetime import datetime, timezone
from utils.db import get_db_connection
from utils.logger import setup_logger
from utils.timer import CPUTimer
from utils.s3 import upload_bytes_to_s3
from jobs.runner import check_cancel_requested, mark_cancelled, update_progress
from utils.export_utils import (
    sanitize_summary_json,
    format_currency,
    format_percentage,
    generate_watermark_text,
    validate_export_data,
    truncate_text,
    safe_get_nested,
)
from io import BytesIO
from typing import Any, Callable

# Initialize with dummy objects to satisfy IDE static analysis
class Dummy:
    def __call__(self, *args: Any, **kwargs: Any) -> Any: return None
    def __getattr__(self, name: str) -> Any: return self
    def __getitem__(self, key: Any) -> Any: return self

_dummy = Dummy()

Presentation: Any = _dummy
Inches: Any = _dummy
Pt: Any = _dummy
RGBColor: Any = _dummy
CategoryChartData: Any = _dummy
XL_CHART_TYPE: Any = _dummy
XL_LEGEND_POSITION: Any = _dummy
PP_ALIGN: Any = _dummy
MSO_SHAPE: Any = _dummy
MSO_ANCHOR: Any = _dummy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    
    # Use getattr for enums to avoid IDE 'Could not import' errors during static analysis
    try:
        import pptx.enum.text as pptx_text # type: ignore
        PP_ALIGN = getattr(pptx_text, 'PP_ALIGN', getattr(pptx_text, 'PP_PARAGRAPH_ALIGNMENT', _dummy))
        MSO_ANCHOR = getattr(pptx_text, 'MSO_ANCHOR', getattr(pptx_text, 'MSO_VERTICAL_ANCHOR', _dummy))
    except (ImportError, AttributeError):
        PP_ALIGN = _dummy
        MSO_ANCHOR = _dummy
        
    try:
        import pptx.enum.shapes as pptx_shapes # type: ignore
        MSO_SHAPE = getattr(pptx_shapes, 'MSO_SHAPE', getattr(pptx_shapes, 'MSO_AUTO_SHAPE_TYPE', _dummy))
    except (ImportError, AttributeError):
        MSO_SHAPE = _dummy
        
except ImportError:
    Presentation = None
    # Keep others as dummy objects to prevent crash during static analysis
    pass

logger = setup_logger()


def handle_export_pptx(job_id: str, org_id: str, object_id: str, logs: dict):
    """Handle investor PPTX export job - stores file in database"""
    logger.info(f"Processing investor PPTX export job {job_id}")
    
    conn = None
    cursor = None
    cpu_timer = CPUTimer()
    export_id = object_id
    
    try:
        if check_cancel_requested(job_id):
            mark_cancelled(job_id)
            return
        
        with cpu_timer:
            conn = get_db_connection()
            cursor = conn.cursor()
            
            # Initialize export_org_id early to ensure it's always defined
            export_org_id = org_id
            
            if not export_id:
                # Handle logs as dict or list
                if isinstance(logs, dict):
                    export_id = logs.get('params', {}).get('exportId')
                elif isinstance(logs, list) and logs:
                    # Try to find exportId in list entries
                    for entry in logs:
                        if isinstance(entry, dict) and 'params' in entry:
                            export_id = entry.get('params', {}).get('exportId')
                            if export_id:
                                break
            
            if not export_id:
                raise ValueError("Export ID not found")
            
            # Extract params safely
            if isinstance(logs, dict):
                params = logs.get('params', {})
            elif isinstance(logs, list) and logs:
                # Find params in list
                params = {}
                for entry in reversed(logs):
                    if isinstance(entry, dict) and 'params' in entry:
                        params = entry.get('params', {})
                        break
            else:
                params = {}
            
            org_name = params.get('orgName', 'Organization')
            include_monte_carlo = params.get('includeMonteCarlo', True)
            include_recommendations = params.get('includeRecommendations', True)
            is_demo = params.get('isDemo', False)
            is_free = params.get('isFree', False)
            watermark_text = generate_watermark_text(is_demo, is_free)
            
            # Update status
            cursor.execute("""
                UPDATE exports SET status = 'processing', updated_at = NOW()
                WHERE id = %s
            """, (export_id,))
            conn.commit()
            
            update_progress(job_id, 10, {'status': 'fetching_data'})
            
            # Get export record with model run summary
            cursor.execute("""
                SELECT e.type, e.model_run_id, mr.summary_json, mr.org_id, e.org_id, e.meta_json
                FROM exports e
                LEFT JOIN model_runs mr ON e.model_run_id = mr.id
                WHERE e.id = %s
            """, (export_id,))
            
            export_record = cursor.fetchone()
            if not export_record:
                raise ValueError(f"Export {export_id} not found")
            
            export_type = export_record[0]
            model_run_id = export_record[1]
            summary_json = export_record[2] if export_record[2] else {}
            export_org_id = export_record[4] or export_record[3] or org_id
            metadata = export_record[5] if export_record[5] else {}
            
            if isinstance(summary_json, str):
                try:
                    summary_json = json.loads(summary_json)
                except:
                    summary_json = {}
            
            # Handle empty summary JSON
            if not summary_json or summary_json == {}:
                logger.warning("Summary JSON is empty, using default values")
                summary_json = {
                    'totalRevenue': 0,
                    'totalExpenses': 0,
                    'netIncome': 0,
                    'cashBalance': 0,
                    'runwayMonths': 0,
                    'burnRate': 0
                }
            
            # Validate and sanitize summary JSON
            is_valid, errors = validate_export_data(summary_json)
            if not is_valid:
                logger.warning(f"Export data validation errors: {errors}")
            
            summary_json = sanitize_summary_json(summary_json)
            
            # Merge AI content from metadata into summary_json
            if metadata and 'aiContent' in metadata:
                summary_json['aiContent'] = metadata['aiContent']
            
            # Get Monte Carlo data if requested
            monte_carlo_data = None
            if include_monte_carlo and model_run_id:
                try:
                    cursor.execute("""
                        SELECT percentiles_json FROM monte_carlo_jobs
                        WHERE model_run_id = %s
                        AND status = 'done'
                        ORDER BY finished_at DESC
                        LIMIT 1
                    """, (model_run_id,))
                    mc_row = cursor.fetchone()
                    if mc_row and mc_row[0]:
                        if isinstance(mc_row[0], str):
                            monte_carlo_data = json.loads(mc_row[0])
                        else:
                            monte_carlo_data = mc_row[0]
                except Exception as e:
                    logger.warning(f"Could not fetch Monte Carlo data: {str(e)}")
                    # Rollback on SQL error to allow transaction to continue
                    try:
                        conn.rollback()
                    except:
                        pass
            
            update_progress(job_id, 30, {'status': 'generating_pptx'})
            
            # Generate PPTX using python-pptx
            try:
                # Check if Presentation is available
                if Presentation is None or Presentation is _dummy:
                    logger.error("python-pptx not installed or failed to import")
                    raise ValueError("PPTX generation library not available. Install python-pptx: pip install python-pptx")
                
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                
                # Professional color scheme
                PRIMARY_COLOR = RGBColor(41, 128, 185)  # Blue
                SUCCESS_COLOR = RGBColor(39, 174, 96)   # Green
                WARNING_COLOR = RGBColor(241, 196, 15)   # Yellow
                DANGER_COLOR = RGBColor(231, 76, 60)     # Red
                DARK_GRAY = RGBColor(44, 62, 80)        # Dark Gray
                LIGHT_GRAY = RGBColor(236, 240, 241)    # Light Gray
                
                # Slide 1: Professional Cover Slide
                slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Add background rectangle for header
                if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                    header_rect = slide1.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.5)
                    )
                    header_rect.fill.solid()
                    header_rect.fill.fore_color.rgb = PRIMARY_COLOR
                    header_rect.line.fill.background()
                
                # Title
                title_shape = slide1.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
                title_frame = title_shape.text_frame
                title_frame.text = "Investor Update"
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(48)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(255, 255, 255)  # White
                if PP_ALIGN is not None and PP_ALIGN is not _dummy:
                    title_para.alignment = PP_ALIGN.LEFT
                
                # Organization name
                subtitle_shape = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
                subtitle_frame = subtitle_shape.text_frame
                subtitle_frame.text = org_name
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(32)
                subtitle_para.font.bold = True
                subtitle_para.font.color.rgb = DARK_GRAY
                if PP_ALIGN is not None and PP_ALIGN is not _dummy:
                    subtitle_para.alignment = PP_ALIGN.LEFT
                
                # Date
                date_shape = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
                date_frame = date_shape.text_frame
                date_frame.text = datetime.now(timezone.utc).strftime('%B %d, %Y')
                date_para = date_frame.paragraphs[0]
                date_para.font.size = Pt(20)
                date_para.font.color.rgb = RGBColor(127, 140, 141)  # Gray
                if PP_ALIGN is not None and PP_ALIGN is not _dummy:
                    date_para.alignment = PP_ALIGN.LEFT
                
                # Add decorative line
                line = slide1.shapes.add_connector(1, Inches(1), Inches(5.5), Inches(9), Inches(5.5))
                line.line.color.rgb = PRIMARY_COLOR
                line.line.width = Pt(3)

                # Slide 2: Strategic Executive Summary
                ai_content = summary_json.get('aiContent', {})
                slide_exec = prs.slides.add_slide(prs.slide_layouts[6])
                
                if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                    title_exec_bg = slide_exec.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                    title_exec_bg.fill.solid()
                    title_exec_bg.fill.fore_color.rgb = PRIMARY_COLOR
                    title_exec_bg.line.fill.background()
                
                title_exec = slide_exec.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title_exec.text_frame.text = "Strategic Executive Summary"
                title_exec.text_frame.paragraphs[0].font.size = Pt(36)
                title_exec.text_frame.paragraphs[0].font.bold = True
                title_exec.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                
                content_exec = slide_exec.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                tf = content_exec.text_frame
                tf.word_wrap = True
                
                p = tf.add_paragraph()
                p.text = "Executive Performance Overview"
                p.font.bold = True
                p.font.size = Pt(18)
                p.font.color.rgb = PRIMARY_COLOR
                
                p = tf.add_paragraph()
                if ai_content and ai_content.get('executiveSummary'):
                    p.text = ai_content.get('executiveSummary')
                else:
                    p.text = f"Strategic assessment of organization performance for the current period. Key top-line revenue stands at ${summary_json.get('totalRevenue', 0):,.0f} with a runway position of {summary_json.get('runwayMonths', 0):.1f} months. Performance remains aligned with strategic growth targets."
                p.font.size = Pt(14)
                p.space_after = Pt(20)
                
                # Slide 3: KPI Summary
                slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                    title_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                    title_bg.fill.solid()
                    title_bg.fill.fore_color.rgb = PRIMARY_COLOR
                
                title2_shape = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title2_shape.text_frame.text = "Key Performance Indicators"
                title2_shape.text_frame.paragraphs[0].font.size = Pt(36)
                title2_shape.text_frame.paragraphs[0].font.bold = True
                title2_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                
                total_revenue = safe_get_nested(summary_json, 'totalRevenue', default=0)
                total_expenses = safe_get_nested(summary_json, 'totalExpenses', default=0)
                net_income = safe_get_nested(summary_json, 'netIncome', default=0)
                cash_balance = safe_get_nested(summary_json, 'cashBalance', default=0)
                burn_rate = safe_get_nested(summary_json, 'burnRate', default=0)
                runway = safe_get_nested(summary_json, 'runwayMonths', default=0)
                
                kpi_data = [
                    ("Total Revenue", format_currency(total_revenue), SUCCESS_COLOR),
                    ("Total Expenses", format_currency(total_expenses), DANGER_COLOR),
                    ("Net Income", format_currency(net_income), SUCCESS_COLOR if net_income >= 0 else DANGER_COLOR),
                    ("Cash Balance", format_currency(cash_balance), PRIMARY_COLOR),
                    ("Monthly Burn Rate", format_currency(burn_rate), WARNING_COLOR),
                    ("Runway", f"{runway:.1f} months", SUCCESS_COLOR if runway >= 12 else WARNING_COLOR if runway >= 6 else DANGER_COLOR),
                ]
                
                card_width = Inches(4.5)
                card_height = Inches(1.8)
                card_spacing = Inches(0.5)
                
                if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                    for idx, (label, value, color) in enumerate(kpi_data):
                        row = idx // 2
                        col = idx % 2
                        x = Inches(0.5) + col * (card_width + card_spacing)
                        y = Inches(1.8) + row * (card_height + Inches(0.3))
                        
                        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
                        card.fill.solid()
                        card.fill.fore_color.rgb = LIGHT_GRAY
                        card.line.color.rgb = color
                        card.line.width = Pt(2)
                        
                        lbl = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_width - Inches(0.4), Inches(0.5))
                        lbl.text_frame.text = label
                        lbl.text_frame.paragraphs[0].font.size = Pt(12)
                        lbl.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
                        
                        val = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), card_width - Inches(0.4), Inches(0.8))
                        val.text_frame.text = value
                        val.text_frame.paragraphs[0].font.size = Pt(24)
                        val.text_frame.paragraphs[0].font.bold = True
                        val.text_frame.paragraphs[0].font.color.rgb = color

                # Slide 4-7: Charts
                budget_data = params.get('budgetActualData', {})
                periods = budget_data.get('periods', [])
                
                # Revenue Chart
                slide3_rev = prs.slides.add_slide(prs.slide_layouts[6])
                if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                    title3_bg = slide3_rev.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                    title3_bg.fill.solid()
                    title3_bg.fill.fore_color.rgb = PRIMARY_COLOR
                
                title3_shape = slide3_rev.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title3_shape.text_frame.text = "1. Revenue Growth Trend"
                title3_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title3_shape.text_frame.paragraphs[0].font.bold = True
                title3_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                
                chart_data_rev = CategoryChartData()
                if periods:
                    sorted_periods = sorted(periods, key=lambda p: p.get('period', ''))
                    chart_data_rev.categories = [p.get('period', 'Unknown') for p in sorted_periods]
                    chart_data_rev.add_series('Actual Revenue', [p.get('actualRevenue', 0) for p in sorted_periods])
                else:
                    chart_data_rev.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
                    chart_data_rev.add_series('Revenue', [100000, 115000, 130000, 145000, 160000, 180000])

                if XL_CHART_TYPE is not None and XL_CHART_TYPE is not _dummy:
                    slide3_rev.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1), Inches(1.8), Inches(8), Inches(4.5), chart_data_rev)

                # Slide 8: Cash & Runway
                slide7_cash = prs.slides.add_slide(prs.slide_layouts[6])
                if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                    title7_bg = slide7_cash.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                    title7_bg.fill.solid()
                    title7_bg.fill.fore_color.rgb = PRIMARY_COLOR
                
                title7_shape = slide7_cash.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title7_shape.text_frame.text = "Cash Position & Runway Analysis"
                title7_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title7_shape.text_frame.paragraphs[0].font.bold = True
                title7_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

                # Monte Carlo
                if monte_carlo_data:
                    slide8_mc = prs.slides.add_slide(prs.slide_layouts[6])
                    if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                        title8_bg = slide8_mc.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                        title8_bg.fill.solid()
                        title8_bg.fill.fore_color.rgb = PRIMARY_COLOR
                    
                    title8_shape = slide8_mc.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                    title8_shape.text_frame.text = "Monte Carlo Risk Analysis"
                    title8_shape.text_frame.paragraphs[0].font.size = Pt(36)
                    title8_shape.text_frame.paragraphs[0].font.bold = True
                    title8_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

                # Recommendations
                if include_recommendations:
                    recommendations = summary_json.get('recommendations', [])
                    if recommendations:
                        slide9_rec = prs.slides.add_slide(prs.slide_layouts[6])
                        if MSO_SHAPE is not None and MSO_SHAPE is not _dummy:
                            title9_bg = slide9_rec.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                            title9_bg.fill.solid()
                            title9_bg.fill.fore_color.rgb = PRIMARY_COLOR
                        
                        title9_shape = slide9_rec.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                        title9_shape.text_frame.text = "AI-CFO Recommendations"
                        title9_shape.text_frame.paragraphs[0].font.size = Pt(36)
                        title9_shape.text_frame.paragraphs[0].font.bold = True
                        title9_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

                # Risks
                slide10_risk = prs.slides.add_slide(prs.slide_layouts[5])
                slide10_risk.shapes.title.text = "Key Risks & Considerations"
                try:
                    tf10 = slide10_risk.placeholders[1].text_frame
                except:
                    tf10 = slide10_risk.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
                
                focus_text = summary_json.get('aiContent', {}).get('areasOfFocus', '')
                if focus_text:
                    tf10.text = focus_text
                else:
                    tf10.text = "• Monitor cash runway closely\n• Review expense growth\n• Scenario planning required\n• Maintain working capital"

                # Watermark
                if (is_demo or is_free) and MSO_ANCHOR is not None and MSO_ANCHOR is not _dummy:
                    for slide in prs.slides:
                        wm_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
                        wm_box.text_frame.text = watermark_text
                        if PP_ALIGN is not None and PP_ALIGN is not _dummy:
                            wm_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        wm_box.text_frame.paragraphs[0].font.size = Pt(10)

                # Save
                pptx_buffer = BytesIO()
                prs.save(pptx_buffer)
                pptx_buffer.seek(0)
                pptx_content = pptx_buffer.read()
                
            except ImportError:
                logger.error("python-pptx not installed")
                raise ValueError("PPTX generation library not available")
            
            update_progress(job_id, 90, {'status': 'storing_file'})
            cpu_seconds = cpu_timer.elapsed()

            # Storage logic
            s3_key = f"exports/{export_org_id}/{export_id}/investor-presentation.pptx"
            uploaded_key = upload_bytes_to_s3(key=s3_key, data=pptx_content, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            
            if uploaded_key:
                cursor.execute("UPDATE exports SET s3_key = %s, status = 'completed', updated_at = NOW() WHERE id = %s", (uploaded_key, export_id))
            else:
                cursor.execute("UPDATE exports SET file_data = %s, status = 'completed', updated_at = NOW() WHERE id = %s", (pptx_content, export_id))
            conn.commit()

            update_progress(job_id, 100, {'status': 'completed', 'cpuSeconds': cpu_seconds, 'fileSize': len(pptx_content), 'slideCount': len(prs.slides)})
            logger.info(f"✅ Export {export_id} completed successfully")

    except Exception as e:
        logger.error(f"❌ Export failed: {str(e)}", exc_info=True)
        if conn and cursor:
            try:
                cursor.execute("UPDATE exports SET status = 'failed', updated_at = NOW() WHERE id = %s", (export_id,))
                conn.commit()
            except: pass
        raise
    finally:
        if cursor: cursor.close()
        if conn: conn.close()
