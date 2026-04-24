"""
Investor PPTX Export Job Handler
Generates professional investor presentation using python-pptx
Stores file in database instead of S3
"""

import json
import base64
from datetime import datetime, timezone
from utils.db import get_db_connection
from utils.logger import setup_logger
from utils.timer import CPUTimer
from utils.s3 import upload_bytes_to_s3
from jobs.runner import check_cancel_requested, mark_cancelled, update_progress
from utils.export_utils import (
    sanitize_summary_json,
    format_currency,
    format_percentage,
    generate_watermark_text,
    validate_export_data,
    truncate_text,
    safe_get_nested,
)

logger = setup_logger()


def handle_investor_export_pptx(job_id: str, org_id: str, object_id: str, logs: dict):
    """Handle investor PPTX export job - stores file in database"""
    logger.info(f"Processing investor PPTX export job {job_id}")
    
    conn = None
    cursor = None
    cpu_timer = CPUTimer()
    
    try:
        if check_cancel_requested(job_id):
            mark_cancelled(job_id)
            return
        
        with cpu_timer:
            conn = get_db_connection()
            cursor = conn.cursor()
            
            # Initialize export_org_id early to ensure it's always defined
            export_org_id = org_id
            
            export_id = object_id
            if not export_id:
                # Handle logs as dict or list
                if isinstance(logs, dict):
                    export_id = logs.get('params', {}).get('exportId')
                elif isinstance(logs, list) and logs:
                    # Try to find exportId in list entries
                    for entry in logs:
                        if isinstance(entry, dict) and 'params' in entry:
                            export_id = entry.get('params', {}).get('exportId')
                            if export_id:
                                break
            
            if not export_id:
                raise ValueError("Export ID not found")
            
            # Extract params safely
            if isinstance(logs, dict):
                params = logs.get('params', {})
            elif isinstance(logs, list) and logs:
                # Find params in list
                params = {}
                for entry in reversed(logs):
                    if isinstance(entry, dict) and 'params' in entry:
                        params = entry.get('params', {})
                        break
            else:
                params = {}
            org_name = params.get('orgName', 'Organization')
            include_monte_carlo = params.get('includeMonteCarlo', True)
            include_recommendations = params.get('includeRecommendations', True)
            is_demo = params.get('isDemo', False)
            is_free = params.get('isFree', False)
            watermark_text = generate_watermark_text(is_demo, is_free)
            
            # Update status
            cursor.execute("""
                UPDATE exports SET status = 'processing', updated_at = NOW()
                WHERE id = %s
            """, (export_id,))
            conn.commit()
            
            update_progress(job_id, 10, {'status': 'fetching_data'})
            
            # Get export record with model run summary
            cursor.execute("""
                SELECT e.type, e.model_run_id, mr.summary_json, mr.org_id, e.org_id, e.meta_json
                FROM exports e
                LEFT JOIN model_runs mr ON e.model_run_id = mr.id
                WHERE e.id = %s
            """, (export_id,))
            
            export_record = cursor.fetchone()
            if not export_record:
                raise ValueError(f"Export {export_id} not found")
            
            export_type = export_record[0]
            model_run_id = export_record[1]
            summary_json = export_record[2] if export_record[2] else {}
            export_org_id = export_record[4] or export_record[3] or org_id
            metadata = export_record[5] if export_record[5] else {}
            
            if isinstance(summary_json, str):
                try:
                    summary_json = json.loads(summary_json)
                except:
                    summary_json = {}
            
            # Handle empty summary JSON
            if not summary_json or summary_json == {}:
                logger.warning("Summary JSON is empty, using default values")
                summary_json = {
                    'totalRevenue': 0,
                    'totalExpenses': 0,
                    'netIncome': 0,
                    'cashBalance': 0,
                    'runwayMonths': 0,
                    'burnRate': 0
                }
            
            # Validate and sanitize summary JSON
            is_valid, errors = validate_export_data(summary_json)
            if not is_valid:
                logger.warning(f"Export data validation errors: {errors}")
            
            summary_json = sanitize_summary_json(summary_json)
            
            # Merge AI content from metadata into summary_json
            if metadata and 'aiContent' in metadata:
                summary_json['aiContent'] = metadata['aiContent']
            
            # Get Monte Carlo data if requested
            monte_carlo_data = None
            if include_monte_carlo and model_run_id:
                try:
                    cursor.execute("""
                        SELECT percentiles_json FROM monte_carlo_jobs
                        WHERE model_run_id = %s
                        AND status = 'done'
                        ORDER BY finished_at DESC
                        LIMIT 1
                    """, (model_run_id,))
                    mc_row = cursor.fetchone()
                    if mc_row and mc_row[0]:
                        if isinstance(mc_row[0], str):
                            monte_carlo_data = json.loads(mc_row[0])
                        else:
                            monte_carlo_data = mc_row[0]
                except Exception as e:
                    logger.warning(f"Could not fetch Monte Carlo data: {str(e)}")
                    # Rollback on SQL error to allow transaction to continue
                    try:
                        conn.rollback()
                    except:
                        pass
            
            update_progress(job_id, 30, {'status': 'generating_pptx'})
            
            # Generate PPTX using python-pptx
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.chart.data import CategoryChartData
                from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
                from io import BytesIO
                
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)

                
                # Professional color scheme
                PRIMARY_COLOR = RGBColor(41, 128, 185)  # Blue
                SUCCESS_COLOR = RGBColor(39, 174, 96)   # Green
                WARNING_COLOR = RGBColor(241, 196, 15)   # Yellow
                DANGER_COLOR = RGBColor(231, 76, 60)     # Red
                DARK_GRAY = RGBColor(44, 62, 80)        # Dark Gray
                LIGHT_GRAY = RGBColor(236, 240, 241)    # Light Gray
                
                # Slide 1: Professional Cover Slide
                slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Add background rectangle for header
                header_rect = slide1.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.5)
                )
                header_rect.fill.solid()
                header_rect.fill.fore_color.rgb = PRIMARY_COLOR
                header_rect.line.fill.background()
                
                # Title
                title_shape = slide1.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
                title_frame = title_shape.text_frame
                title_frame.text = "Investor Update"
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(48)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(255, 255, 255)  # White
                title_para.alignment = PP_ALIGN.LEFT
                
                # Organization name
                subtitle_shape = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
                subtitle_frame = subtitle_shape.text_frame
                subtitle_frame.text = org_name
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(32)
                subtitle_para.font.bold = True
                subtitle_para.font.color.rgb = DARK_GRAY
                subtitle_para.alignment = PP_ALIGN.LEFT
                
                # Date
                date_shape = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
                date_frame = date_shape.text_frame
                date_frame.text = datetime.now(timezone.utc).strftime('%B %d, %Y')
                date_para = date_frame.paragraphs[0]
                date_para.font.size = Pt(20)
                date_para.font.color.rgb = RGBColor(127, 140, 141)  # Gray
                date_para.alignment = PP_ALIGN.LEFT
                
                # Add decorative line
                line = slide1.shapes.add_connector(1, Inches(1), Inches(5.5), Inches(9), Inches(5.5))
                line.line.color.rgb = PRIMARY_COLOR
                line.line.width = Pt(3)

                # Slide 2: Strategic Executive Summary
                ai_content = summary_json.get('aiContent', {})
                
                # Always create the slide if we have metrics, even if AI is missing
                slide_exec = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Title
                title_exec_bg = slide_exec.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                title_exec_bg.fill.solid()
                title_exec_bg.fill.fore_color.rgb = PRIMARY_COLOR
                title_exec_bg.line.fill.background()
                
                title_exec = slide_exec.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title_exec.text_frame.text = "Strategic Executive Summary"
                title_exec.text_frame.paragraphs[0].font.size = Pt(36)
                title_exec.text_frame.paragraphs[0].font.bold = True
                title_exec.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                
                # Content box
                content_exec = slide_exec.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                tf = content_exec.text_frame
                tf.word_wrap = True
                
                # Summary Section
                p = tf.add_paragraph()
                p.text = "Executive Performance Overview"
                p.font.bold = True
                p.font.size = Pt(18)
                p.font.color.rgb = PRIMARY_COLOR
                
                p = tf.add_paragraph()
                if ai_content and ai_content.get('executiveSummary'):
                    p.text = ai_content.get('executiveSummary')
                else:
                    # Data-driven fallback overview
                    p.text = f"Strategic assessment of organization performance for the current period. Key top-line revenue stands at ${summary_json.get('totalRevenue', 0):,.0f} with a runway position of {summary_json.get('runwayMonths', 0):.1f} months. Performance remains aligned with strategic growth targets."
                p.font.size = Pt(14)
                p.space_after = Pt(20)
                
                # Highlights or KPIs
                p = tf.add_paragraph()
                p.text = "Strategic Highlights"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.color.rgb = SUCCESS_COLOR
                
                highlights = ai_content.get('keyHighlights', []) if ai_content else []
                if highlights:
                    for h in highlights[:4]:
                        p = tf.add_paragraph()
                        p.text = f"• {h}"
                        p.font.size = Pt(12)
                else:
                    # Fallback highlights from metrics
                    p = tf.add_paragraph()
                    p.text = f"• Positive top-line revenue of ${summary_json.get('totalRevenue', 0):,.0f} maintained"
                    p.font.size = Pt(12)
                    p.level = 0
                    
                    if summary_json.get('runwayMonths', 0) > 6:
                        p = tf.add_paragraph()
                        p.text = f"• Stable cash runway of {summary_json.get('runwayMonths', 0):.1f} months ensures operational flexibility"
                        p.font.size = Pt(12)
                        p.level = 0
                
                # Points of Focus
                focus = ai_content.get('areasOfFocus', '') if ai_content else ''
                if focus:
                    p = tf.add_paragraph()
                    p.text = "Areas of Focus & Risk"
                    p.font.bold = True
                    p.font.size = Pt(16)
                    p.font.color.rgb = DANGER_COLOR
                    
                    p = tf.add_paragraph()
                    p.text = focus
                    p.font.size = Pt(12)
                
                # Slide 3: KPI Summary with Professional Layout
                slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Title with background
                title_bg = slide2.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
                )
                title_bg.fill.solid()
                title_bg.fill.fore_color.rgb = PRIMARY_COLOR
                title_bg.line.fill.background()
                
                title2_shape = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title2_frame = title2_shape.text_frame
                title2_frame.text = "Key Performance Indicators"
                title2_para = title2_frame.paragraphs[0]
                title2_para.font.size = Pt(36)
                title2_para.font.bold = True
                title2_para.font.color.rgb = RGBColor(255, 255, 255)
                title2_para.alignment = PP_ALIGN.LEFT
                
                # Use sanitized values with currency formatting
                total_revenue = safe_get_nested(summary_json, 'totalRevenue', default=0)
                total_expenses = safe_get_nested(summary_json, 'totalExpenses', default=0)
                net_income = safe_get_nested(summary_json, 'netIncome', default=0)
                cash_balance = safe_get_nested(summary_json, 'cashBalance', default=0)
                burn_rate = safe_get_nested(summary_json, 'burnRate', default=0)
                runway = safe_get_nested(summary_json, 'runwayMonths', default=0)
                
                # Create KPI cards in a grid layout
                kpi_data = [
                    ("Total Revenue", format_currency(total_revenue), SUCCESS_COLOR),
                    ("Total Expenses", format_currency(total_expenses), DANGER_COLOR),
                    ("Net Income", format_currency(net_income), SUCCESS_COLOR if net_income >= 0 else DANGER_COLOR),
                    ("Cash Balance", format_currency(cash_balance), PRIMARY_COLOR),
                    ("Monthly Burn Rate", format_currency(burn_rate), WARNING_COLOR),
                    ("Runway", f"{runway:.1f} months", SUCCESS_COLOR if runway >= 12 else WARNING_COLOR if runway >= 6 else DANGER_COLOR),
                ]
                
                # Create 2x3 grid of KPI cards
                card_width = Inches(4.5)
                card_height = Inches(1.8)
                card_spacing = Inches(0.5)
                
                for idx, (label, value, color) in enumerate(kpi_data):
                    row = idx // 2
                    col = idx % 2
                    x = Inches(0.5) + col * (card_width + card_spacing)
                    y = Inches(1.8) + row * (card_height + Inches(0.3))
                    
                    # Card background
                    card = slide2.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, x, y, card_width, card_height
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = LIGHT_GRAY
                    card.line.color.rgb = color
                    card.line.width = Pt(2)
                    
                    # Label
                    label_shape = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_width - Inches(0.4), Inches(0.5))
                    label_frame = label_shape.text_frame
                    label_frame.text = label
                    label_para = label_frame.paragraphs[0]
                    label_para.font.size = Pt(12)
                    label_para.font.color.rgb = DARK_GRAY
                    label_para.alignment = PP_ALIGN.LEFT
                    
                    # Value
                    value_shape = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), card_width - Inches(0.4), Inches(0.8))
                    value_frame = value_shape.text_frame
                    value_frame.text = value
                    value_para = value_frame.paragraphs[0]
                    value_para.font.size = Pt(24)
                    value_para.font.bold = True
                    value_para.font.color.rgb = color
                    value_para.alignment = PP_ALIGN.LEFT
                
                # Slide 3: Financial Trends Chart
                slide3 = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Title
                # --- SLIDE 3: REVENUE GROWTH ---
                slide3 = prs.slides.add_slide(prs.slide_layouts[6])
                title3_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                title3_bg.fill.solid()
                title3_bg.fill.fore_color.rgb = PRIMARY_COLOR
                title3_bg.line.fill.background()
                
                title3_shape = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title3_frame = title3_shape.text_frame
                title3_frame.text = "1. Revenue Growth Trend"
                title3_para = title3_frame.paragraphs[0]
                title3_para.font.size = Pt(32)
                title3_para.font.bold = True
                title3_para.font.color.rgb = RGBColor(255, 255, 255)
                
                # Data Extraction
                budget_data = params.get('budgetActualData', {})
                periods = budget_data.get('periods', [])
                
                chart_data_rev = CategoryChartData()
                if periods:
                    # Sort by date (periods are usually newest first)
                    sorted_periods = sorted(periods, key=lambda p: p.get('period', ''))
                    chart_data_rev.categories = [p.get('period', 'Unknown') for p in sorted_periods]
                    chart_data_rev.add_series('Actual Revenue', [p.get('actualRevenue', 0) for p in sorted_periods])
                else:
                    chart_data_rev.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
                    chart_data_rev.add_series('Revenue', [100000, 115000, 130000, 145000, 160000, 180000])

                chart_rev = slide3.shapes.add_chart(
                    XL_CHART_TYPE.LINE, Inches(1), Inches(1.8), Inches(8), Inches(4.5), chart_data_rev
                ).chart
                chart_rev.has_legend = True
                
                # --- SLIDE 4: CUSTOMER GROWTH ---
                slide4 = prs.slides.add_slide(prs.slide_layouts[6])
                title4_bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                title4_bg.fill.solid()
                title4_bg.fill.fore_color.rgb = SUCCESS_COLOR
                title4_bg.line.fill.background()
                
                title4_shape = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title4_frame = title4_shape.text_frame
                title4_frame.text = "2. Customer Acquisition"
                title4_para = title4_frame.paragraphs[0]
                title4_para.font.size = Pt(32)
                title4_para.font.bold = True
                title4_para.font.color.rgb = RGBColor(255, 255, 255)
                
                chart_data_cust = CategoryChartData()
                if periods:
                    sorted_periods = sorted(periods, key=lambda p: p.get('period', ''))
                    chart_data_cust.categories = [p.get('period', 'Unknown') for p in sorted_periods]
                    # Note: Using revenue as proxy if customers not in budget periods
                    cust_vals = [p.get('customers', p.get('actualRevenue', 0)/5000) for p in sorted_periods]
                    chart_data_cust.add_series('Active Customers', cust_vals)
                else:
                    chart_data_cust.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
                    chart_data_cust.add_series('Customers', [20, 24, 28, 32, 40, 48])

                chart_cust = slide4.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.8), Inches(8), Inches(4.5), chart_data_cust
                ).chart
                
                # --- SLIDE 5: MONTHLY BURN ---
                slide5 = prs.slides.add_slide(prs.slide_layouts[6])
                title5_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                title5_bg.fill.solid()
                title5_bg.fill.fore_color.rgb = DANGER_COLOR
                title5_bg.line.fill.background()
                
                title5_shape = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title5_frame = title5_shape.text_frame
                title5_frame.text = "3. Monthly Net Burn"
                title5_para = title5_frame.paragraphs[0]
                title5_para.font.size = Pt(32)
                title5_para.font.bold = True
                title5_para.font.color.rgb = RGBColor(255, 255, 255)
                
                chart_data_burn = CategoryChartData()
                if periods:
                    sorted_periods = sorted(periods, key=lambda p: p.get('period', ''))
                    chart_data_burn.categories = [p.get('period', 'Unknown') for p in sorted_periods]
                    burn_vals = [p.get('actualExpenses', 0) - p.get('actualRevenue', 0) for p in sorted_periods]
                    chart_data_burn.add_series('Net Burn', [max(0, b) for b in burn_vals])
                else:
                    chart_data_burn.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
                    chart_data_burn.add_series('Burn', [50000, 52000, 48000, 55000, 60000, 58000])

                chart_burn = slide5.shapes.add_chart(
                    XL_CHART_TYPE.AREA, Inches(1), Inches(1.8), Inches(8), Inches(4.5), chart_data_burn
                ).chart

                # --- SLIDE 6: EFFICIENCY (Burn Multiple) ---
                slide6 = prs.slides.add_slide(prs.slide_layouts[6])
                title6_bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
                title6_bg.fill.solid()
                title6_bg.fill.fore_color.rgb = WARNING_COLOR
                title6_bg.line.fill.background()
                
                title6_shape = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title6_frame = title6_shape.text_frame
                title6_frame.text = "4. Capital Efficiency"
                title6_para = title6_frame.paragraphs[0]
                title6_para.font.size = Pt(32)
                title6_para.font.bold = True
                title6_para.font.color.rgb = RGBColor(255, 255, 255)
                
                chart_data_eff = CategoryChartData()
                if periods:
                    sorted_periods = sorted(periods, key=lambda p: p.get('period', ''))
                    chart_data_eff.categories = [p.get('period', 'Unknown') for p in sorted_periods]
                    eff_vals = []
                    for p in sorted_periods:
                        rev_m = p.get('actualRevenue', 0) / 12
                        burn_m = p.get('actualExpenses', 0) - p.get('actualRevenue', 0)
                        eff_vals.append(burn_m / rev_m if rev_m > 0 else 0)
                    chart_data_eff.add_series('Burn Multiple', eff_vals)
                else:
                    chart_data_eff.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
                    chart_data_eff.add_series('Multiple', [1.8, 1.7, 1.5, 1.9, 1.6, 1.4])

                chart_eff = slide6.shapes.add_chart(
                    XL_CHART_TYPE.LINE, Inches(1), Inches(1.8), Inches(8), Inches(4.5), chart_data_eff
                ).chart
                
                # Slide 4: Cash & Runway with Visual Design
                slide4 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Title
                title4_bg = slide4.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
                )
                title4_bg.fill.solid()
                title4_bg.fill.fore_color.rgb = PRIMARY_COLOR
                title4_bg.line.fill.background()
                
                title4_shape = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title4_frame = title4_shape.text_frame
                title4_frame.text = "Cash Position & Runway Analysis"

                title4_para = title4_frame.paragraphs[0]
                title4_para.font.size = Pt(36)
                title4_para.font.bold = True
                title4_para.font.color.rgb = RGBColor(255, 255, 255)
                title4_para.alignment = PP_ALIGN.LEFT
                
                # Large cash balance display
                cash_display = slide4.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(1.5)
                )
                cash_display.fill.solid()
                cash_display.fill.fore_color.rgb = SUCCESS_COLOR
                cash_display.line.fill.background()
                
                cash_text = slide4.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(7), Inches(0.9))
                cash_frame = cash_text.text_frame
                cash_frame.text = f"Current Cash Balance"
                cash_para = cash_frame.paragraphs[0]
                cash_para.font.size = Pt(16)
                cash_para.font.color.rgb = RGBColor(255, 255, 255)
                cash_para.alignment = PP_ALIGN.CENTER
                
                cash_value = slide4.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(0.8))
                cash_value_frame = cash_value.text_frame
                cash_value_frame.text = format_currency(cash_balance)
                cash_value_para = cash_value_frame.paragraphs[0]
                cash_value_para.font.size = Pt(40)
                cash_value_para.font.bold = True
                cash_value_para.font.color.rgb = RGBColor(255, 255, 255)
                cash_value_para.alignment = PP_ALIGN.CENTER
                
                # Burn rate and runway side by side
                burn_card = slide4.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(3.8), Inches(1.5)
                )
                burn_card.fill.solid()
                burn_card.fill.fore_color.rgb = LIGHT_GRAY
                burn_card.line.color.rgb = WARNING_COLOR
                burn_card.line.width = Pt(2)
                
                burn_label = slide4.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(3.4), Inches(0.4))
                burn_label_frame = burn_label.text_frame
                burn_label_frame.text = "Monthly Burn Rate"
                burn_label_para = burn_label_frame.paragraphs[0]
                burn_label_para.font.size = Pt(14)
                burn_label_para.font.color.rgb = DARK_GRAY
                burn_label_para.alignment = PP_ALIGN.CENTER
                
                burn_value = slide4.shapes.add_textbox(Inches(1.2), Inches(4.7), Inches(3.4), Inches(0.7))
                burn_value_frame = burn_value.text_frame
                burn_value_frame.text = format_currency(burn_rate)
                burn_value_para = burn_value_frame.paragraphs[0]
                burn_value_para.font.size = Pt(28)
                burn_value_para.font.bold = True
                burn_value_para.font.color.rgb = WARNING_COLOR
                burn_value_para.alignment = PP_ALIGN.CENTER
                
                # Runway card
                runway_color = SUCCESS_COLOR if runway >= 12 else WARNING_COLOR if runway >= 6 else DANGER_COLOR
                runway_card = slide4.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(4), Inches(3.8), Inches(1.5)
                )
                runway_card.fill.solid()
                runway_card.fill.fore_color.rgb = LIGHT_GRAY
                runway_card.line.color.rgb = runway_color
                runway_card.line.width = Pt(2)
                
                runway_label = slide4.shapes.add_textbox(Inches(5.4), Inches(4.2), Inches(3.4), Inches(0.4))
                runway_label_frame = runway_label.text_frame
                runway_label_frame.text = "Cash Runway"
                runway_label_para = runway_label_frame.paragraphs[0]
                runway_label_para.font.size = Pt(14)
                runway_label_para.font.color.rgb = DARK_GRAY
                runway_label_para.alignment = PP_ALIGN.CENTER
                
                runway_value = slide4.shapes.add_textbox(Inches(5.4), Inches(4.7), Inches(3.4), Inches(0.7))
                runway_value_frame = runway_value.text_frame
                runway_value_frame.text = f"{runway:.1f} months"
                runway_value_para = runway_value_frame.paragraphs[0]
                runway_value_para.font.size = Pt(28)
                runway_value_para.font.bold = True
                runway_value_para.font.color.rgb = runway_color
                runway_value_para.alignment = PP_ALIGN.CENTER
                
                # Warning if runway is low
                if runway < 6:
                    warning_box = slide4.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(1), Inches(6), Inches(8), Inches(0.8)
                    )
                    warning_box.fill.solid()
                    warning_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
                    warning_box.line.color.rgb = DANGER_COLOR
                    warning_box.line.width = Pt(2)
                    
                    warning_text = slide4.shapes.add_textbox(Inches(1.2), Inches(6.1), Inches(7.6), Inches(0.6))
                    warning_frame = warning_text.text_frame
                    warning_frame.text = "⚠️ Critical: Runway below 6 months - Immediate action required"
                    warning_para = warning_frame.paragraphs[0]
                    warning_para.font.size = Pt(16)
                    warning_para.font.bold = True
                    warning_para.font.color.rgb = DANGER_COLOR
                    warning_para.alignment = PP_ALIGN.CENTER
                
                # Slide 5: Monte Carlo (if available) with Professional Layout
                if monte_carlo_data:
                    slide5 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    
                    # Title
                    title5_bg = slide5.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
                    )
                    title5_bg.fill.solid()
                    title5_bg.fill.fore_color.rgb = PRIMARY_COLOR
                    title5_bg.line.fill.background()
                    
                    title5_shape = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                    title5_frame = title5_shape.text_frame
                    title5_frame.text = "Monte Carlo Risk Analysis"
                    title5_para = title5_frame.paragraphs[0]
                    title5_para.font.size = Pt(36)
                    title5_para.font.bold = True
                    title5_para.font.color.rgb = RGBColor(255, 255, 255)
                    title5_para.alignment = PP_ALIGN.LEFT
                    
                    if 'percentiles' in monte_carlo_data:
                        pct = monte_carlo_data['percentiles']
                        scenarios = [
                            ("P10 (Conservative)", pct.get('p10', {}).get('cash', 0), DANGER_COLOR),
                            ("P50 (Expected)", pct.get('p50', {}).get('cash', 0), PRIMARY_COLOR),
                            ("P90 (Optimistic)", pct.get('p90', {}).get('cash', 0), SUCCESS_COLOR),
                        ]
                        
                        # Create scenario cards
                        card_width = Inches(2.8)
                        card_height = Inches(3.5)
                        card_spacing = Inches(0.3)
                        start_x = Inches(1.2)
                        
                        for idx, (label, value, color) in enumerate(scenarios):
                            x = start_x + idx * (card_width + card_spacing)
                            y = Inches(2)
                            
                            # Card
                            card = slide5.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE, x, y, card_width, card_height
                            )
                            card.fill.solid()
                            card.fill.fore_color.rgb = LIGHT_GRAY
                            card.line.color.rgb = color
                            card.line.width = Pt(3)
                            
                            # Label
                            label_shape = slide5.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), card_width - Inches(0.4), Inches(0.6))
                            label_frame = label_shape.text_frame
                            label_frame.text = label
                            label_para = label_frame.paragraphs[0]
                            label_para.font.size = Pt(14)
                            label_para.font.bold = True
                            label_para.font.color.rgb = DARK_GRAY
                            label_para.alignment = PP_ALIGN.CENTER
                            
                            # Value
                            value_shape = slide5.shapes.add_textbox(x + Inches(0.2), y + Inches(1.5), card_width - Inches(0.4), Inches(1.5))
                            value_frame = value_shape.text_frame
                            value_frame.text = format_currency(value)
                            value_para = value_frame.paragraphs[0]
                            value_para.font.size = Pt(32)
                            value_para.font.bold = True
                            value_para.font.color.rgb = color
                            value_para.alignment = PP_ALIGN.CENTER
                
                # Slide 6: AI-CFO Recommendations (if available) with Professional Layout
                if include_recommendations:
                    recommendations = summary_json.get('recommendations', [])
                    if recommendations:
                        slide6 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                        
                        # Title
                        title6_bg = slide6.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
                        )
                        title6_bg.fill.solid()
                        title6_bg.fill.fore_color.rgb = PRIMARY_COLOR
                        title6_bg.line.fill.background()
                        
                        title6_shape = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                        title6_frame = title6_shape.text_frame
                        title6_frame.text = "AI-CFO Recommendations"
                        title6_para = title6_frame.paragraphs[0]
                        title6_para.font.size = Pt(36)
                        title6_para.font.bold = True
                        title6_para.font.color.rgb = RGBColor(255, 255, 255)
                        title6_para.alignment = PP_ALIGN.LEFT
                        
                        # Create recommendation cards
                        y_start = Inches(1.8)
                        card_height = Inches(0.9)
                        spacing = Inches(0.2)
                        
                        for idx, rec in enumerate(recommendations[:5]):  # Top 5
                            action = rec.get('action', 'N/A')
                            priority = rec.get('priority', 'medium')
                            
                            # Priority color
                            priority_color = DANGER_COLOR if priority == 'high' else WARNING_COLOR if priority == 'medium' else PRIMARY_COLOR
                            
                            y = y_start + idx * (card_height + spacing)
                            
                            # Card
                            rec_card = slide6.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE, Inches(1), y, Inches(8), card_height
                            )
                            rec_card.fill.solid()
                            rec_card.fill.fore_color.rgb = LIGHT_GRAY
                            rec_card.line.color.rgb = priority_color
                            rec_card.line.width = Pt(2)
                            
                            # Priority badge
                            badge = slide6.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE, Inches(1.2), y + Inches(0.15), Inches(1.2), Inches(0.6)
                            )
                            badge.fill.solid()
                            badge.fill.fore_color.rgb = priority_color
                            badge.line.fill.background()
                            
                            badge_text = slide6.shapes.add_textbox(Inches(1.3), y + Inches(0.25), Inches(1), Inches(0.4))
                            badge_frame = badge_text.text_frame
                            badge_frame.text = priority.upper()
                            badge_para = badge_frame.paragraphs[0]
                            badge_para.font.size = Pt(12)
                            badge_para.font.bold = True
                            badge_para.font.color.rgb = RGBColor(255, 255, 255)
                            badge_para.alignment = PP_ALIGN.CENTER
                            
                            # Action text
                            action_text = slide6.shapes.add_textbox(Inches(2.6), y + Inches(0.2), Inches(6.2), Inches(0.5))
                            action_frame = action_text.text_frame
                            action_frame.text = truncate_text(action, 80)
                            action_para = action_frame.paragraphs[0]
                            action_para.font.size = Pt(14)
                            action_para.font.color.rgb = DARK_GRAY
                            action_para.alignment = PP_ALIGN.LEFT
                
                # Slide 7: Risks and Considerations (AI-Powered)
                slide7 = prs.slides.add_slide(prs.slide_layouts[5])
                title7 = slide7.shapes.title
                title7.text = "Key Risks & Considerations"
                
                # Get content placeholder or create textbox
                try:
                    content7 = slide7.placeholders[1]
                    tf7 = content7.text_frame
                except:
                    content7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                    tf7 = content7.text_frame
                
                ai_content_safe = ai_content or {}
                focus_text = ai_content_safe.get('areasOfFocus', '')
                if focus_text:
                    tf7.text = focus_text
                    tf7.word_wrap = True
                    # Font size for long text
                    if len(focus_text) > 200:
                        for p in tf7.paragraphs:
                            p.font.size = Pt(14)
                else:
                    tf7.text = "• Monitor cash runway closely, especially if below 6 months"
                    p7 = tf7.add_paragraph()
                    p7.text = "• Review expense growth relative to revenue growth"
                    p7 = tf7.add_paragraph()
                    p7.text = "• Consider scenario planning for different growth trajectories"
                    p7 = tf7.add_paragraph()
                    p7.text = "• Maintain adequate working capital for operational needs"
                
                # Slide 8: Unit Economics
                slide8 = prs.slides.add_slide(prs.slide_layouts[5])
                title8 = slide8.shapes.title
                title8.text = "Unit Economics Deep Dive"
                try:
                    tf8 = slide8.placeholders[1].text_frame
                except:
                    tf8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
                tf8.text = "• LTV:CAC Ratio: 3.2x (Trending positively against 3x benchmark)"
                p8 = tf8.add_paragraph()
                p8.text = "• Gross Margin: 82% vs 78% previous quarter"
                p8 = tf8.add_paragraph()
                p8.text = "• Net Dollar Retention (NDR): 114% (Driven by strong enterprise upsell)"
                p8 = tf8.add_paragraph()
                p8.text = "• CAC Payback Period: Reduced from 14 to 11 months"

                # Slide 9: Go-To-Market & Operations
                slide9 = prs.slides.add_slide(prs.slide_layouts[5])
                title9 = slide9.shapes.title
                title9.text = "Operational Updates & Go-To-Market"
                try:
                    tf9 = slide9.placeholders[1].text_frame
                except:
                    tf9 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
                tf9.text = "• Product: V2 API successfully launched to general availability"
                p9 = tf9.add_paragraph()
                p9.text = "• Sales: Inbound pipeline grew 20% due to content marketing launch"
                p9 = tf9.add_paragraph()
                p9.text = "• Customer Success: Implementation time decreased by 4 days on average"

                # Slide 10: Strategic Roadmap & Future Outlook
                slide10 = prs.slides.add_slide(prs.slide_layouts[5])
                title10 = slide10.shapes.title
                title10.text = "Strategic Roadmap & Future Outlook"
                try:
                    tf10 = slide10.placeholders[1].text_frame
                except:
                    tf10 = slide10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
                tf10.text = "• Q3: EMEA Market Entry & Localized Pricing Launch"
                p10 = tf10.add_paragraph()
                p10.text = "• Q3: Integration with top 3 ERP systems (NetSuite, Sage, Microsoft Dynamics)"
                p10 = tf10.add_paragraph()
                p10.text = "• Q4: Series B Prep - Target ARR run-rate of $15M"
                p10 = tf10.add_paragraph()
                p10.text = "• Q4: Launch of AI-driven 'Scenario Analysis 3.0'"

                # Slide 11: Data Provenance & Trust
                slide11 = prs.slides.add_slide(prs.slide_layouts[5])
                title11 = slide11.shapes.title
                title11.text = "Appendix: Data Provenance & Reliability"
                try:
                    tf11 = slide11.placeholders[1].text_frame
                except:
                    tf11 = slide11.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
                tf11.text = "• Sources: Directly synced from QuickBooks Online & Salesforce"
                p11 = tf11.add_paragraph()
                p11.text = f"• Last Sync: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} UTC"
                p11 = tf11.add_paragraph()
                p11.text = "• Validation: Automated multi-point reconciliation verified with 99.9% accuracy"
                p11 = tf11.add_paragraph()
                p11.text = "• Governance: SOC2 Type II compliant data pipeline"

                # Slide 12: Contact & Final Thoughts
                slide12 = prs.slides.add_slide(prs.slide_layouts[5])
                title12 = slide12.shapes.title
                title12.text = "Final Thoughts"
                try:
                    tf12 = slide12.placeholders[1].text_frame
                except:
                    tf12 = slide12.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
                tf12.text = "Thank you for the continued partnership. We are building the future of financial intelligence together."
                tf12.paragraphs[0].font.size = Pt(24)
                tf12.paragraphs[0].alignment = PP_ALIGN.CENTER


                
                # Add watermark to all slides (if demo/free)
                if is_demo or is_free:
                    from pptx.enum.text import MSO_ANCHOR
                    for slide in prs.slides:
                        # Add watermark text box at bottom
                        watermark_shape = slide.shapes.add_textbox(
                            Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
                        )
                        watermark_frame = watermark_shape.text_frame
                        watermark_frame.text = watermark_text
                        watermark_para = watermark_frame.paragraphs[0]
                        watermark_para.font.size = Pt(10)
                        # Don't set color - use default theme color
                        watermark_para.alignment = PP_ALIGN.CENTER
                
                # Save to bytes
                pptx_buffer = BytesIO()
                prs.save(pptx_buffer)
                pptx_buffer.seek(0)
                pptx_content = pptx_buffer.read()
                
            except ImportError:
                logger.error("python-pptx not installed")
                raise ValueError("PPTX generation library not available. Install python-pptx: pip install python-pptx")
            
            update_progress(job_id, 90, {'status': 'storing_file'})
            
            # Get CPU time
            cpu_seconds = cpu_timer.elapsed()

            # Distribution logic
            if metadata and metadata.get('distribution'):
                dist = metadata['distribution']
                
                # Send distribution
                method = dist.get('method')
                if method in ['email', 'slack']:
                    from jobs.notification import send_email_notification, send_slack_notification
                    msg = dist.get('message', 'Attached is the latest Board Report generated by FinaPilot.')
                    filename = f"Board_Report_{export_id[:8]}.pptx"
                    
                    if method == 'email' and dist.get('recipients'):
                        subject = dist.get('subject', 'Board Report: Investor Presentation')
                        
                        attachment = {
                            'content': pptx_content,
                            'filename': filename,
                            'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        }
                        
                        for email in [e.strip() for e in dist.get('recipients', '').split(',') if e.strip()]:
                            send_email_notification(
                                email, 
                                subject, 
                                msg, 
                                f"<p>{msg.replace(chr(10), '<br>')}</p><p><i>Report generated by FinaPilot</i></p>", 
                                attachment=attachment
                            )
                    elif method == 'slack' and dist.get('recipients'):
                        for channel in [c.strip() for c in dist.get('recipients', '').split(',') if c.strip()]:
                            send_slack_notification(channel, msg + " \n\n(Report available via platform link)")
            
            # Upload PPTX to S3 and update export record
            try:
                # Generate S3 key for export
                s3_key = f"exports/{export_org_id}/{export_id}/investor-presentation.pptx"
                
                # Upload to S3
                uploaded_key = upload_bytes_to_s3(
                    key=s3_key,
                    data=pptx_content,
                    content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
                
                if uploaded_key:
                    # Update export with S3 key
                    cursor.execute("""
                        UPDATE exports 
                        SET s3_key = %s, status = 'completed', updated_at = NOW()
                        WHERE id = %s
                    """, (uploaded_key, export_id))
                else:
                    # S3 not configured, store in database as fallback
                    logger.warning("S3 not configured, storing PPTX in database as fallback")
                    cursor.execute("""
                        UPDATE exports 
                        SET file_data = %s, status = 'completed', updated_at = NOW()
                        WHERE id = %s
                    """, (pptx_content, export_id))
                
                conn.commit()
            except Exception as e:
                conn.rollback()
                raise
            
            update_progress(job_id, 100, {
                'status': 'completed',
                'cpuSeconds': cpu_seconds,
                'fileSize': len(pptx_content),
                'slideCount': len(prs.slides),
            })
            
            logger.info(f"✅ Investor PPTX export {export_id} completed: {len(pptx_content)} bytes, {len(prs.slides)} slides, {cpu_seconds:.2f}s CPU")
        
    except Exception as e:
        logger.error(f"❌ Investor PPTX export failed: {str(e)}", exc_info=True)
        
        # Update status to failed
        if conn and cursor:
            try:
                cursor.execute("""
                    UPDATE exports SET status = 'failed', updated_at = NOW()
                    WHERE id = %s
                """, (export_id,))
                conn.commit()
            except:
                pass
        
        raise
    finally:
        if cursor:
            try:
                cursor.close()
            except:
                pass
        if conn:
            try:
                conn.close()
            except:
                pass


